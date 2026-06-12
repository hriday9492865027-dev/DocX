# AetherPDF Backend - Python Document Transmutation Server
import os
import uuid
import shutil
import sys
import subprocess
from flask import Flask, request, send_file, jsonify
from flask_cors import CORS

# Import libraries for local sandboxed conversions
from pdf2docx import Converter
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import io
import pdfplumber
import openpyxl

app = Flask(__name__)
CORS(app)

TEMP_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'temp')
os.makedirs(TEMP_DIR, exist_ok=True)

# Conditional imports for Windows COM automation
IS_WINDOWS = sys.platform.startswith('win')
if IS_WINDOWS:
    try:
        import pythoncom
        import win32com.client
    except ImportError:
        print("win32com or pythoncom not available. Running in LibreOffice-only mode.")
        IS_WINDOWS = False

# ==========================================
# --- CROSS-PLATFORM CONVERSION HELPERS ---
# ==========================================

def convert_with_libreoffice(input_path, output_dir):
    """Converts a document to PDF using LibreOffice headless command line."""
    # Find soffice executable path
    soffice_path = shutil.which('soffice') or shutil.which('libreoffice')
    
    if not soffice_path and sys.platform.startswith('win'):
        # Check standard Windows directories if not in PATH
        possible_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"
        ]
        for p in possible_paths:
            if os.path.exists(p):
                soffice_path = p
                break
                
    if not soffice_path:
        raise Exception(
            "LibreOffice (soffice) CLI not found. Make sure LibreOffice is installed "
            "and added to system variables/PATH (or MS Office is installed on Windows)."
        )
        
    cmd = [
        soffice_path,
        '--headless',
        '--convert-to', 'pdf',
        input_path,
        '--outdir', output_dir
    ]
    
    print(f"Executing LibreOffice conversion: {' '.join(cmd)}")
    result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    if result.returncode != 0:
        raise Exception(f"LibreOffice conversion failed: {result.stderr or result.stdout}")

# ==========================================
# --- CONVERSION ALGORITHMS ---
# ==========================================

# 1. PDF to DOCX
def pdf_to_docx(pdf_path, docx_path):
    cv = Converter(pdf_path)
    cv.convert(docx_path, start=0, end=None)
    cv.close()

def get_transparent_image_bytes(doc, xref, smask_xref):
    try:
        base_pix = fitz.Pixmap(doc, xref)
        
        # If it's CMYK and has no smask, convert to RGB
        if base_pix.colorspace and base_pix.colorspace.n == 4 and smask_xref == 0:
            rgb_pix = fitz.Pixmap(fitz.csRGB, base_pix)
            img_bytes = rgb_pix.tobytes("png")
            return img_bytes
            
        if smask_xref > 0:
            mask_pix = fitz.Pixmap(doc, smask_xref)
            # Make sure base is RGB before combining (fitz.Pixmap expects RGB base for masking)
            if base_pix.colorspace and base_pix.colorspace.n != 3:
                base_pix = fitz.Pixmap(fitz.csRGB, base_pix)
            
            pix = fitz.Pixmap(base_pix, mask_pix)
            img_bytes = pix.tobytes("png")
            return img_bytes
        else:
            # Check colorspace to ensure compatibility with python-pptx (RGB)
            if base_pix.colorspace and base_pix.colorspace.n != 3:
                base_pix = fitz.Pixmap(fitz.csRGB, base_pix)
            img_bytes = base_pix.tobytes("png")
            return img_bytes
    except Exception as e:
        print(f"Error extracting transparent image (xref={xref}, smask={smask_xref}): {e}")
        return None

# 2. PDF to PPTX (Generates editable text elements and extracts images)
def pdf_to_pptx(pdf_path, pptx_path):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    
    doc = fitz.open(pdf_path)
    prs = Presentation()
    blank_layout = prs.slide_layouts[6] # blank layout
    
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        
        # Set slide size to match page size (72 points = 1 inch)
        if page_num == 0:
            prs.slide_width = Inches(page.rect.width / 72.0)
            prs.slide_height = Inches(page.rect.height / 72.0)
            
        slide = prs.slides.add_slide(blank_layout)
        page_dict = page.get_text("dict")
        img_list = page.get_images(full=True)
        used_img_indices = set()
        
        # Separate blocks to enforce Z-order: background images first, then foreground text
        text_blocks = []
        image_blocks = []
        
        for block in page_dict.get("blocks", []):
            block_type = block.get("type", 0)
            if block_type == 0:
                text_blocks.append(block)
            elif block_type == 1:
                image_blocks.append(block)
                
        # 1. Render all image blocks first (Background Layer)
        for block in image_blocks:
            bbox = block.get("bbox", (0, 0, 0, 0))
            x0, y0, x1, y1 = bbox
            
            left = Inches(x0 / 72.0)
            top = Inches(y0 / 72.0)
            width = Inches((x1 - x0) / 72.0)
            height = Inches((y1 - y0) / 72.0)
            
            image_bytes = block.get("image", None)
            
            # Match image by aspect ratio to extract transparency if available
            xref = None
            smask_xref = 0
            block_width = block.get("width", 1)
            block_height = block.get("height", 1)
            block_aspect = block_width / block_height if block_height > 0 else 1
            
            # Find the closest matching image in the page's image catalog
            best_idx = -1
            min_diff = 999999.0
            for idx, img in enumerate(img_list):
                if idx in used_img_indices:
                    continue
                img_width = img[2]
                img_height = img[3]
                if img_height > 0:
                    img_aspect = img_width / img_height
                    diff = abs(block_aspect - img_aspect)
                    if diff < min_diff:
                        min_diff = diff
                        best_idx = idx
                        
            # Only match if the aspect ratios are reasonably close (within 15% tolerance)
            if best_idx != -1 and min_diff < 0.15:
                img = img_list[best_idx]
                xref = img[0]
                smask_xref = img[1]
                used_img_indices.add(best_idx)
                
            if xref is not None:
                transparent_bytes = get_transparent_image_bytes(doc, xref, smask_xref)
                if transparent_bytes:
                    image_bytes = transparent_bytes
                    
            if image_bytes:
                try:
                    slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, width, height)
                except Exception as img_err:
                    print(f"Failed to insert PDF image block: {img_err}")
                    
        # 2. Render all text blocks second (Foreground Layer)
        for block in text_blocks:
            bbox = block.get("bbox", (0, 0, 0, 0))
            x0, y0, x1, y1 = bbox
            
            left = Inches(x0 / 72.0)
            top = Inches(y0 / 72.0)
            width = Inches((x1 - x0) / 72.0)
            height = Inches((y1 - y0) / 72.0)
            
            lines = block.get("lines", [])
            if not lines:
                continue
                
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            
            tf.margin_left = Inches(0.02)
            tf.margin_right = Inches(0.02)
            tf.margin_top = Inches(0.02)
            tf.margin_bottom = Inches(0.02)
            
            for i, line in enumerate(lines):
                spans = line.get("spans", [])
                if not spans:
                    continue
                    
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                    
                for span in spans:
                    span_text = span.get("text", "")
                    if not span_text:
                        continue
                        
                    run = p.add_run()
                    run.text = span_text
                    run.font.size = Pt(span.get("size", 12))
                    
                    color_int = span.get("color", 0)
                    r = (color_int >> 16) & 255
                    g = (color_int >> 8) & 255
                    b = color_int & 255
                    run.font.color.rgb = RGBColor(r, g, b)
                    
                    font_name = span.get("font", "").lower()
                    if "bold" in font_name:
                        run.font.bold = True
                    if "italic" in font_name:
                        run.font.italic = True
                        
    prs.save(pptx_path)

# 3. PDF to XLSX
def pdf_to_xlsx(pdf_path, xlsx_path):
    wb = openpyxl.Workbook()
    wb.remove(wb.active) # Remove default sheet
    
    with pdfplumber.open(pdf_path) as pdf:
        for i, page in enumerate(pdf.pages):
            sheet = wb.create_sheet(title=f"Page {i+1}")
            tables = page.extract_tables()
            row_idx = 1
            
            for table in tables:
                for row in table:
                    for col_idx, cell in enumerate(row, start=1):
                        sheet.cell(row=row_idx, column=col_idx, value=cell)
                    row_idx += 1
                row_idx += 2 # spacing between tables
                
            if not tables:
                text = page.extract_text()
                if text:
                    for line in text.split('\n'):
                        parts = [p.strip() for p in line.split('   ') if p.strip()]
                        if not parts:
                            parts = line.split('\t')
                        for col_idx, val in enumerate(parts, start=1):
                            sheet.cell(row=row_idx, column=col_idx, value=val)
                        row_idx += 1
                        
    if not wb.sheetnames:
        wb.create_sheet(title="Sheet 1")
    wb.save(xlsx_path)

# 4. DOCX to PDF
def docx_to_pdf(docx_path, pdf_path):
    if IS_WINDOWS:
        try:
            pythoncom.CoInitialize()
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            doc = word.Documents.Open(os.path.abspath(docx_path))
            doc.SaveAs(os.path.abspath(pdf_path), FileFormat=17) # 17 is wdFormatPDF
            doc.Close()
            word.Quit()
            return
        except Exception as e:
            print(f"Windows COM automation failed: {e}. Falling back to LibreOffice CLI...")
        finally:
            pythoncom.CoUninitialize()
            
    # Linux / Fallback: Run LibreOffice
    output_dir = os.path.dirname(pdf_path)
    convert_with_libreoffice(docx_path, output_dir)
    
    filename_no_ext = os.path.basename(docx_path).rsplit('.', 1)[0]
    expected_out = os.path.join(output_dir, f"{filename_no_ext}.pdf")
    if os.path.exists(expected_out):
        shutil.move(expected_out, pdf_path)
    else:
        raise Exception("LibreOffice CLI did not produce PDF output.")

# 5. PPTX to PDF
def pptx_to_pdf(pptx_path, pdf_path):
    if IS_WINDOWS:
        try:
            pythoncom.CoInitialize()
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            pres = powerpoint.Presentations.Open(os.path.abspath(pptx_path), WithWindow=False)
            pres.SaveAs(os.path.abspath(pdf_path), FileFormat=32) # 32 is ppSaveAsPDF
            pres.Close()
            powerpoint.Quit()
            return
        except Exception as e:
            print(f"Windows COM automation failed: {e}. Falling back to LibreOffice CLI...")
        finally:
            pythoncom.CoUninitialize()
            
    # Linux / Fallback: Run LibreOffice
    output_dir = os.path.dirname(pdf_path)
    convert_with_libreoffice(pptx_path, output_dir)
    
    filename_no_ext = os.path.basename(pptx_path).rsplit('.', 1)[0]
    expected_out = os.path.join(output_dir, f"{filename_no_ext}.pdf")
    if os.path.exists(expected_out):
        shutil.move(expected_out, pdf_path)
    else:
        raise Exception("LibreOffice CLI did not produce PDF output.")

# 6. XLSX to PDF
def xlsx_to_pdf(xlsx_path, pdf_path):
    if IS_WINDOWS:
        try:
            pythoncom.CoInitialize()
            excel = win32com.client.Dispatch("Excel.Application")
            excel.Visible = False
            wb = excel.Workbooks.Open(os.path.abspath(xlsx_path))
            wb.ExportAsFixedFormat(0, os.path.abspath(pdf_path)) # 0 is xlTypePDF
            wb.Close(SaveChanges=False)
            excel.Quit()
            return
        except Exception as e:
            print(f"Windows COM automation failed: {e}. Falling back to LibreOffice CLI...")
        finally:
            pythoncom.CoUninitialize()
            
    # Linux / Fallback: Run LibreOffice
    output_dir = os.path.dirname(pdf_path)
    convert_with_libreoffice(xlsx_path, output_dir)
    
    filename_no_ext = os.path.basename(xlsx_path).rsplit('.', 1)[0]
    expected_out = os.path.join(output_dir, f"{filename_no_ext}.pdf")
    if os.path.exists(expected_out):
        shutil.move(expected_out, pdf_path)
    else:
        raise Exception("LibreOffice CLI did not produce PDF output.")

# ==========================================
# --- ROUTE HANDLERS ---
# ==========================================

@app.route('/convert', methods=['POST'])
def convert_file():
    if 'file' not in request.files or 'mode' not in request.form:
        return jsonify({"error": "Missing file or mode parameters"}), 400
        
    uploaded_file = request.files['file']
    mode = request.form['mode']
    
    if uploaded_file.filename == '':
        return jsonify({"error": "No file selected"}), 400
        
    # Generate unique filenames to prevent conflicts
    task_id = str(uuid.uuid4())
    filename = uploaded_file.filename
    ext = filename.split('.')[-1].lower()
    
    src_path = os.path.join(TEMP_DIR, f"{task_id}.{ext}")
    dest_ext = 'pdf' if 'to-pdf' in mode else ('docx' if 'to-doc' in mode else ('pptx' if 'to-ppt' in mode else 'xlsx'))
    dest_path = os.path.join(TEMP_DIR, f"{task_id}_out.{dest_ext}")
    
    # Save input file
    uploaded_file.save(src_path)
    
    try:
        if mode == 'pdf-to-doc':
            pdf_to_docx(src_path, dest_path)
        elif mode == 'pdf-to-ppt':
            pdf_to_pptx(src_path, dest_path)
        elif mode == 'pdf-to-sheets':
            pdf_to_xlsx(src_path, dest_path)
        elif mode == 'doc-to-pdf':
            docx_to_pdf(src_path, dest_path)
        elif mode == 'ppt-to-pdf':
            pptx_to_pdf(src_path, dest_path)
        elif mode == 'sheets-to-pdf':
            xlsx_to_pdf(src_path, dest_path)
        else:
            return jsonify({"error": f"Invalid conversion mode: {mode}"}), 400
            
        if not os.path.exists(dest_path):
            raise Exception("Output file was not generated by conversion module.")
            
        # Return converted file
        return send_file(
            dest_path,
            as_attachment=True,
            download_name=f"{filename.rsplit('.', 1)[0]}_converted.{dest_ext}"
        )
        
    except Exception as e:
        print(f"Error during conversion: {str(e)}")
        return jsonify({"error": str(e)}), 500
        
    finally:
        # Cleanup temp files
        try:
            if os.path.exists(src_path):
                os.remove(src_path)
            if os.path.exists(dest_path):
                os.remove(dest_path)
        except Exception as cleanup_err:
            print(f"Failed to delete temp files: {cleanup_err}")

if __name__ == '__main__':
    # Binding to 0.0.0.0 is necessary for running inside Docker
    port = int(os.environ.get("PORT", 5000))
    print(f"AetherPDF Python backend listening on http://0.0.0.0:{port}...")
    app.run(host="0.0.0.0", port=port, debug=False)
