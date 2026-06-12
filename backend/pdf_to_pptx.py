import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import io

def get_transparent_image_bytes(doc, xref, smask_xref):
    try:
        base_pix = fitz.Pixmap(doc, xref)
        
        # If it's CMYK and has no smask, convert to RGB
        if base_pix.colorspace and base_pix.colorspace.n == 4 and smask_xref == 0:
            rgb_pix = fitz.Pixmap(fitz.csRGB, base_pix)
            img_bytes = rgb_pix.tobytes("png")
            return img_bytes
            
        if smask_xref > 0:
            mask_pix = fitz.Pixmap(doc, smask_xref)
            # Make sure base is RGB before combining (fitz.Pixmap expects RGB base for masking)
            if base_pix.colorspace and base_pix.colorspace.n != 3:
                base_pix = fitz.Pixmap(fitz.csRGB, base_pix)
            
            pix = fitz.Pixmap(base_pix, mask_pix)
            img_bytes = pix.tobytes("png")
            return img_bytes
        else:
            # Check colorspace to ensure compatibility with python-pptx (RGB)
            if base_pix.colorspace and base_pix.colorspace.n != 3:
                base_pix = fitz.Pixmap(fitz.csRGB, base_pix)
            img_bytes = base_pix.tobytes("png")
            return img_bytes
    except Exception as e:
        print(f"Error extracting transparent image (xref={xref}, smask={smask_xref}): {e}")
        return None

def pdf_to_pptx(pdf_path, pptx_path):
    doc = fitz.open(pdf_path)
    prs = Presentation()
    blank_layout = prs.slide_layouts[6] # blank layout
    
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        
        # Set slide size to match page size (72 points = 1 inch)
        if page_num == 0:
            prs.slide_width = Inches(page.rect.width / 72.0)
            prs.slide_height = Inches(page.rect.height / 72.0)
            
        slide = prs.slides.add_slide(blank_layout)
        page_dict = page.get_text("dict")
        img_list = page.get_images(full=True)
        used_img_indices = set()
        
        # Separate blocks to enforce Z-order: background images first, then foreground text
        text_blocks = []
        image_blocks = []
        
        for block in page_dict.get("blocks", []):
            block_type = block.get("type", 0)
            if block_type == 0:
                text_blocks.append(block)
            elif block_type == 1:
                image_blocks.append(block)
                
        # 1. Render all image blocks first (Background Layer)
        for block in image_blocks:
            bbox = block.get("bbox", (0, 0, 0, 0))
            x0, y0, x1, y1 = bbox
            
            left = Inches(x0 / 72.0)
            top = Inches(y0 / 72.0)
            width = Inches((x1 - x0) / 72.0)
            height = Inches((y1 - y0) / 72.0)
            
            image_bytes = block.get("image", None)
            
            # Match image by aspect ratio to extract transparency if available
            xref = None
            smask_xref = 0
            block_width = block.get("width", 1)
            block_height = block.get("height", 1)
            block_aspect = block_width / block_height if block_height > 0 else 1
            
            # Find the closest matching image in the page's image catalog
            best_idx = -1
            min_diff = 999999.0
            for idx, img in enumerate(img_list):
                if idx in used_img_indices:
                    continue
                img_width = img[2]
                img_height = img[3]
                if img_height > 0:
                    img_aspect = img_width / img_height
                    diff = abs(block_aspect - img_aspect)
                    if diff < min_diff:
                        min_diff = diff
                        best_idx = idx
                        
            # Only match if the aspect ratios are reasonably close (within 15% tolerance)
            if best_idx != -1 and min_diff < 0.15:
                img = img_list[best_idx]
                xref = img[0]
                smask_xref = img[1]
                used_img_indices.add(best_idx)
                
            if xref is not None and smask_xref > 0:
                transparent_bytes = get_transparent_image_bytes(doc, xref, smask_xref)
                if transparent_bytes:
                    image_bytes = transparent_bytes
                    
            if image_bytes:
                try:
                    slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, width, height)
                except Exception as img_err:
                    print(f"Failed to insert PDF image block: {img_err}")
                    
        # 2. Render all text blocks second (Foreground Layer)
        for block in text_blocks:
            bbox = block.get("bbox", (0, 0, 0, 0))
            x0, y0, x1, y1 = bbox
            
            left = Inches(x0 / 72.0)
            top = Inches(y0 / 72.0)
            width = Inches((x1 - x0) / 72.0)
            height = Inches((y1 - y0) / 72.0)
            
            lines = block.get("lines", [])
            if not lines:
                continue
                
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            
            tf.margin_left = Inches(0.02)
            tf.margin_right = Inches(0.02)
            tf.margin_top = Inches(0.02)
            tf.margin_bottom = Inches(0.02)
            
            for i, line in enumerate(lines):
                spans = line.get("spans", [])
                if not spans:
                    continue
                    
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                    
                for span in spans:
                    span_text = span.get("text", "")
                    if not span_text:
                        continue
                        
                    run = p.add_run()
                    run.text = span_text
                    run.font.size = Pt(span.get("size", 12))
                    
                    color_int = span.get("color", 0)
                    r = (color_int >> 16) & 255
                    g = (color_int >> 8) & 255
                    b = color_int & 255
                    run.font.color.rgb = RGBColor(r, g, b)
                    
                    font_name = span.get("font", "").lower()
                    if "bold" in font_name:
                        run.font.bold = True
                    if "italic" in font_name:
                        run.font.italic = True
                        
    prs.save(pptx_path)
